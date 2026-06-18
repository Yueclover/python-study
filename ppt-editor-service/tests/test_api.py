import io
from pptx import Presentation
from fastapi.testclient import TestClient
import app.main as main_mod
from app.main import app


def _setup_storage(tmp_path):
    from app.storage import Storage
    main_mod.storage = Storage(str(tmp_path))


def test_parse_then_apply_then_download(tmp_path, basic_pptx_path):
    _setup_storage(tmp_path)
    client = TestClient(app)

    with open(basic_pptx_path, "rb") as f:
        resp = client.post("/parse", files={"file": ("t.pptx", f, "application/octet-stream")})
    assert resp.status_code == 200
    doc = resp.json()
    doc_id = doc["doc_id"]
    assert doc["slides"][0]["shapes"][0]["shape_id"] == "s1_sh1"

    resp2 = client.post("/apply", json={
        "doc_id": doc_id,
        "ops": [{"op": "set_text", "shape_id": "s1_sh1", "text": "改后标题"}],
    })
    assert resp2.status_code == 200
    body = resp2.json()
    assert body["applied"] == 1 and body["rejected"] == []
    name = body["download_url"].split("/files/")[1]

    resp3 = client.get(f"/files/{name}")
    assert resp3.status_code == 200
    prs = Presentation(io.BytesIO(resp3.content))
    assert prs.slides[0].shapes[0].text_frame.text == "改后标题"


def test_parse_rejects_garbage(tmp_path):
    _setup_storage(tmp_path)
    client = TestClient(app)
    resp = client.post("/parse", files={"file": ("x.pptx", io.BytesIO(b"not a pptx"), "application/octet-stream")})
    assert resp.status_code == 400


def test_apply_unknown_doc(tmp_path):
    _setup_storage(tmp_path)
    client = TestClient(app)
    resp = client.post("/apply", json={"doc_id": "ghost", "ops": []})
    assert resp.status_code == 404
