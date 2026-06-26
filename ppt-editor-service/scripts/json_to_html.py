"""Demo: 把 parse_presentation 返回的 JSON 转成静态 HTML 预览。

用法:
    python -m scripts.json_to_html storage/104fba29/source.pptx out.html

说明:版式近似 + 文字/表格准确;图片/图表只画占位框(parser 未导出二进制)。
"""
import sys
import html
import base64
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 让 "python scripts/json_to_html.py" 也能 import app
sys.path.insert(0, ".")
from app.parser import parse_presentation  # noqa: E402
from app.ids import shape_id_for  # noqa: E402
from scripts.visuals import (  # noqa: E402
    load_theme, extract_visuals, slide_bg_css,
)

EMU_PER_PX = 9525  # 914400 EMU/inch ÷ 96 px/inch


def extract_images(prs):
    """按 shape_id 抽取图片为 data URI,不污染 parse 主 JSON。

    返回 {shape_id: "data:image/png;base64,..."}。与 parser 同序遍历,
    复用 shape_id_for 保证 id 一致。
    """
    out = {}
    for si, slide in enumerate(prs.slides):
        for oi, shp in enumerate(slide.shapes):
            if shp.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                img = shp.image
                uri = (f"data:{img.content_type};base64,"
                       + base64.b64encode(img.blob).decode())
                out[shape_id_for(si, oi)] = uri
            except Exception:
                pass  # 占位图/损坏图 → 留空,降级为占位框
    return out


def px(emu):
    return 0 if emu is None else round(emu / EMU_PER_PX)


def esc(t):
    return html.escape(t or "").replace("\n", "<br>")


_GEOM_CSS = {
    "oval": "border-radius:50%;",
    "rounded": "border-radius:10px;",
    "triangle": "clip-path:polygon(50% 0,100% 100%,0 100%);",
}


def shape_html(sh, images, visuals):
    p = sh["pos"]
    style = (
        f"position:absolute;left:{px(p['x'])}px;top:{px(p['y'])}px;"
        f"width:{px(p['w'])}px;height:{px(p['h'])}px;"
    )
    vis = visuals.get(sh["shape_id"])
    if vis:
        style += vis.get("transform", "")
        if vis.get("svg"):
            # FREEFORM:用 SVG <path> 重建任意曲线
            return f"<div class='shape' style='{style}'>{vis['svg']}</div>"
        if vis.get("bg"):
            style += f"background:{vis['bg']};"
        if vis.get("border"):
            style += f"border:{vis['border']};"
        if vis.get("geom"):
            style += _GEOM_CSS.get(vis["geom"], "")
    st = sh.get("style") or {}
    if st.get("size"):
        style += f"font-size:{st['size']}px;"
    if st.get("bold"):
        style += "font-weight:bold;"
    if st.get("color"):
        style += f"color:{st['color']};"
    if st.get("align"):
        style += f"text-align:{st['align']};"
    if st.get("font"):
        style += f"font-family:'{st['font']}',sans-serif;"

    kind = sh["type"]
    if kind == "table" and "table" in sh:
        t = sh["table"]
        grid = {(c["r"], c["c"]): c["text"] for c in t["cells"]}
        rows = "".join(
            "<tr>" + "".join(
                f"<td>{esc(grid.get((r, c), ''))}</td>" for c in range(t["cols"])
            ) + "</tr>"
            for r in range(t["rows"])
        )
        inner = f"<table class='tbl'>{rows}</table>"
        return f"<div class='shape' style='{style}'>{inner}</div>"

    if kind == "picture":
        uri = images.get(sh["shape_id"])
        if uri:
            return (
                f"<div class='shape' style='{style}'>"
                f"<img src='{uri}' style='width:100%;height:100%;"
                f"object-fit:contain;display:block'></div>"
            )
        # 抽取失败 → 占位框
        return (
            f"<div class='shape ph' style='{style}'>"
            f"<span class='ph-label'>🖼 图片</span></div>"
        )
    if kind == "chart":
        return (
            f"<div class='shape ph' style='{style}'>"
            f"<span class='ph-label'>📊 图表</span></div>"
        )

    return f"<div class='shape' style='{style}'>{esc(sh.get('text', ''))}</div>"


def to_html(parsed, images=None, visuals=None, backgrounds=None):
    images = images or {}
    visuals = visuals or {}
    backgrounds = backgrounds or {}
    W = px(parsed["slide_size"]["width"])
    H = px(parsed["slide_size"]["height"])
    slides = []
    for sl in parsed["slides"]:
        shapes = "".join(shape_html(sh, images, visuals) for sh in sl["shapes"])
        bg = backgrounds.get(sl["slide_id"])
        bg_css = f";background:{bg}" if bg else ""
        slides.append(
            f"<div class='slide-wrap'><div class='slide-no'>#{sl['index'] + 1} "
            f"· {esc(sl['layout_name'])}</div>"
            f"<div class='slide' style='width:{W}px;height:{H}px{bg_css}'>{shapes}</div></div>"
        )
    body = "\n".join(slides)
    return f"""<!doctype html><html lang="zh"><meta charset="utf-8">
<title>PPT 预览</title>
<style>
  body{{background:#525659;margin:0;padding:24px;font-family:'Microsoft YaHei',sans-serif}}
  .slide-wrap{{margin:0 auto 28px;width:{W}px}}
  .slide-no{{color:#bbb;font-size:13px;margin-bottom:6px}}
  .slide{{position:relative;background:#fff;box-shadow:0 4px 18px rgba(0,0,0,.4);
          overflow:hidden}}
  .shape{{overflow:hidden;box-sizing:border-box;line-height:1.25}}
  .ph{{border:1px dashed #bbb;display:flex;align-items:center;justify-content:center;
       background:#f5f5f5}}
  .ph-label{{color:#999;font-size:14px}}
  .tbl{{border-collapse:collapse;width:100%;height:100%}}
  .tbl td{{border:1px solid #ccc;padding:2px 6px;font-size:12px}}
</style>
{body}
</html>"""


if __name__ == "__main__":
    src = sys.argv[1] if len(sys.argv) > 1 else "storage/104fba29/source.pptx"
    out = sys.argv[2] if len(sys.argv) > 2 else "preview.html"
    prs = Presentation(src)
    parsed = parse_presentation(prs)
    images = extract_images(prs)
    theme = load_theme(prs)
    visuals = extract_visuals(prs, theme)
    backgrounds = {
        s["slide_id"]: slide_bg_css(prs.slides[s["index"]], prs, theme)
        for s in parsed["slides"]
    }
    with open(out, "w", encoding="utf-8") as f:
        f.write(to_html(parsed, images, visuals, backgrounds))
    fills = sum(1 for v in visuals.values() if v["bg"] and not v["is_freeform"])
    print(f"wrote {out}: {len(parsed['slides'])} slides, {len(images)} images, "
          f"{fills} shape-fills, {sum(1 for b in backgrounds.values() if b)} bg")
