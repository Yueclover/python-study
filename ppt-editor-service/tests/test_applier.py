from pptx import Presentation
from app.applier import apply_ops


def test_set_text_applied(basic_pptx_path):
    prs = Presentation(basic_pptx_path)
    applied, rejected = apply_ops(prs, [
        {"op": "set_text", "shape_id": "s1_sh1", "text": "新标题"},
    ])
    assert applied == 1 and rejected == []
    assert prs.slides[0].shapes[0].text_frame.text == "新标题"


def test_dup_then_fill_copies(basic_pptx_path):
    prs = Presentation(basic_pptx_path)
    applied, rejected = apply_ops(prs, [
        {"op": "dup_slide", "slide_id": "s2", "count": 2, "as": ["d1", "d2"]},
        {"op": "set_text", "shape_id": "d1::sh1", "text": "要点一"},
        {"op": "set_text", "shape_id": "d2::sh1", "text": "要点二"},
    ])
    assert applied == 3 and rejected == []
    assert len(prs.slides) == 4
    assert prs.slides[2].shapes[0].text_frame.text == "要点一"
    assert prs.slides[3].shapes[0].text_frame.text == "要点二"


def test_bad_id_rejected_but_others_apply(basic_pptx_path):
    prs = Presentation(basic_pptx_path)
    applied, rejected = apply_ops(prs, [
        {"op": "set_text", "shape_id": "s9_sh9", "text": "x"},
        {"op": "set_text", "shape_id": "s1_sh1", "text": "ok"},
    ])
    assert applied == 1
    assert len(rejected) == 1
    assert rejected[0]["index"] == 0
    assert "不存在" in rejected[0]["reason"]
    assert prs.slides[0].shapes[0].text_frame.text == "ok"


def test_del_slide(basic_pptx_path):
    prs = Presentation(basic_pptx_path)
    applied, rejected = apply_ops(prs, [{"op": "del_slide", "slide_id": "s1"}])
    assert applied == 1 and rejected == []
    assert len(prs.slides) == 1


def test_non_operror_exception_isolated(table_pptx_path):
    """TypeError from bad rows type must be locally rejected; subsequent op must still apply."""
    prs = Presentation(table_pptx_path)
    # s1_sh2 is the table; rows="5" (string) triggers TypeError deep in set_table_size
    # s1_sh1 is the title text shape; set_text on it should still succeed
    ops = [
        {"op": "set_table_size", "shape_id": "s1_sh2", "rows": "5", "cols": 4},
        {"op": "set_text", "shape_id": "s1_sh1", "text": "隔离测试标题"},
    ]
    applied, rejected = apply_ops(prs, ops)

    # Bad op is isolated: exactly one rejection at index 0
    assert len(rejected) == 1, f"expected 1 rejected, got {rejected}"
    assert rejected[0]["index"] == 0
    assert rejected[0]["op"] == "set_table_size"
    assert "执行失败" in rejected[0]["reason"]

    # Valid op still applied
    assert applied == 1
    assert prs.slides[0].shapes[0].text_frame.text == "隔离测试标题"
