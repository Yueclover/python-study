from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN

from .ids import slide_id_for, shape_id_for

_ALIGN = {
    PP_ALIGN.LEFT: "left",
    PP_ALIGN.CENTER: "center",
    PP_ALIGN.RIGHT: "right",
    PP_ALIGN.JUSTIFY: "justify",
}


def _shape_kind(shape):
    if shape.has_table:
        return "table"
    if getattr(shape, "has_chart", False):
        return "chart"
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    if shape.has_text_frame:
        return "text"
    return "other"


def _ph_type(shape):
    if not shape.is_placeholder:
        return None
    t = shape.placeholder_format.type
    if t in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
        return "title"
    if t == PP_PLACEHOLDER.SUBTITLE:
        return "subtitle"
    if t in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
        return "body"
    return str(t).split()[0].lower()


def _pos(shape):
    def emu(v):
        return int(v) if v is not None else None
    return {"x": emu(shape.left), "y": emu(shape.top),
            "w": emu(shape.width), "h": emu(shape.height)}


def _style(shape):
    tf = shape.text_frame
    para = tf.paragraphs[0]
    run = para.runs[0] if para.runs else None
    font = run.font if run else para.font
    color = None
    try:
        if font.color and font.color.type is not None:
            color = "#" + str(font.color.rgb)
    except Exception:
        color = None
    return {
        "font": font.name,
        "size": int(font.size.pt) if font.size else None,
        "bold": bool(font.bold) if font.bold is not None else False,
        "color": color,
        "align": _ALIGN.get(para.alignment),
    }


def _table(shape):
    t = shape.table
    rows = len(t.rows)
    cols = len(t.columns)
    cells = [{"r": r, "c": c, "text": t.cell(r, c).text}
             for r in range(rows) for c in range(cols)]
    return {"rows": rows, "cols": cols, "cells": cells}


def _parse_shape(shape, slide_index0, order0):
    kind = _shape_kind(shape)
    d = {
        "shape_id": shape_id_for(slide_index0, order0),
        "type": kind,
        "name": shape.name,
        "ph_type": _ph_type(shape),
        "pos": _pos(shape),
    }
    if shape.has_text_frame:
        d["text"] = shape.text_frame.text
        d["style"] = _style(shape)
    if kind == "table":
        d["table"] = _table(shape)
    return d


def _parse_slide(slide, index0):
    return {
        "slide_id": slide_id_for(index0),
        "index": index0,
        "layout_name": slide.slide_layout.name,
        "shapes": [_parse_shape(s, index0, i) for i, s in enumerate(slide.shapes)],
    }


def parse_presentation(prs):
    return {
        "slide_size": {"width": int(prs.slide_width), "height": int(prs.slide_height)},
        "slides": [_parse_slide(s, i) for i, s in enumerate(prs.slides)],
    }
