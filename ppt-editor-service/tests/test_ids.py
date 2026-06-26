from pptx import Presentation
from app.ids import slide_id_for, shape_id_for, IdIndex


def test_id_encoding():
    assert slide_id_for(0) == "s1"
    assert shape_id_for(0, 2) == "s1_sh3"


def test_index_resolves_slide_and_shape(basic_pptx_path):
    prs = Presentation(basic_pptx_path)
    idx = IdIndex(prs)
    assert idx.slide("s1") is not None
    assert idx.slide("s99") is None
    # s1 第一个形状是标题占位符
    shp = idx.shape("s1_sh1")
    assert shp is not None
    assert shp.has_text_frame


def test_temp_slide_reference(basic_pptx_path):
    prs = Presentation(basic_pptx_path)
    idx = IdIndex(prs)
    src = idx.slide("s2")
    idx.register_temp_slide("d1", src)
    assert idx.slide("d1") is src
    # d1::sh1 解析为 d1 这页里第 1 个形状
    assert idx.shape("d1::sh1") == list(src.shapes)[0]
