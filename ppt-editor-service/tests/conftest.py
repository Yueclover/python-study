import pytest
from pptx import Presentation
from pptx.util import Inches


@pytest.fixture
def basic_pptx_path(tmp_path):
    """两页：标题页 + 要点页。"""
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[0])  # 标题幻灯片
    s1.shapes.title.text = "原标题"
    s1.placeholders[1].text = "原副标题"
    s2 = prs.slides.add_slide(prs.slide_layouts[1])  # 标题和内容
    s2.shapes.title.text = "要点页标题"
    s2.placeholders[1].text = "要点占位"
    p = tmp_path / "basic.pptx"
    prs.save(str(p))
    return str(p)


@pytest.fixture
def table_pptx_path(tmp_path):
    """一页：仅标题 + 一个 3x4 表格。"""
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])  # 仅标题
    gf = s.shapes.add_table(3, 4, Inches(1), Inches(2), Inches(8), Inches(3))
    gf.table.cell(0, 0).text = "季度"
    p = tmp_path / "table.pptx"
    prs.save(str(p))
    return str(p)
