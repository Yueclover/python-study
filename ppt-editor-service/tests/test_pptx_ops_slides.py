from pptx import Presentation
from app import pptx_ops


def test_duplicate_after_keeps_order_and_content(basic_pptx_path):
    prs = Presentation(basic_pptx_path)
    src = prs.slides[1]  # 要点页
    news = pptx_ops.duplicate_slide_after(prs, src, 2)
    assert len(news) == 2
    assert len(prs.slides) == 4
    # 副本紧跟在源页（index 1）之后
    assert pptx_ops.slide_index(prs, news[0]) == 2
    assert pptx_ops.slide_index(prs, news[1]) == 3
    # 副本内容复制自源页标题
    assert news[0].shapes[0].text_frame.text == "要点页标题"


def test_delete_slide(basic_pptx_path):
    prs = Presentation(basic_pptx_path)
    target = prs.slides[0]
    assert pptx_ops.delete_slide(prs, target) is True
    assert len(prs.slides) == 1
