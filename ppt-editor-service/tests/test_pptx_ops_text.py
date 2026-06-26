from pptx import Presentation
from app import pptx_ops


def test_set_text_keeps_bold(basic_pptx_path):
    prs = Presentation(basic_pptx_path)
    title = prs.slides[0].shapes[0]
    run = title.text_frame.paragraphs[0].runs[0]
    run.font.bold = True
    pptx_ops.set_text_keep_style(title.text_frame, "新标题")
    assert title.text_frame.text == "新标题"
    assert title.text_frame.paragraphs[0].runs[0].font.bold is True


def test_set_cell(table_pptx_path):
    prs = Presentation(table_pptx_path)
    tbl_shape = [s for s in prs.slides[0].shapes if s.has_table][0]
    pptx_ops.set_cell(tbl_shape, 1, 2, "85%")
    assert tbl_shape.table.cell(1, 2).text == "85%"


def test_set_table_size_add_and_remove_rows(table_pptx_path):
    prs = Presentation(table_pptx_path)
    tbl_shape = [s for s in prs.slides[0].shapes if s.has_table][0]
    pptx_ops.set_table_size(tbl_shape, rows=5, cols=4)
    assert len(tbl_shape.table.rows) == 5
    pptx_ops.set_table_size(tbl_shape, rows=2, cols=4)
    assert len(tbl_shape.table.rows) == 2


def test_set_table_size_add_col(table_pptx_path):
    prs = Presentation(table_pptx_path)
    tbl_shape = [s for s in prs.slides[0].shapes if s.has_table][0]
    pptx_ops.set_table_size(tbl_shape, rows=3, cols=6)
    assert len(tbl_shape.table.columns) == 6


def test_add_column_does_not_inherit_last_col_text(table_pptx_path):
    """新增列时，复制自末列的单元格不应继承其文本。"""
    prs = Presentation(table_pptx_path)
    tbl_shape = [s for s in prs.slides[0].shapes if s.has_table][0]
    table = tbl_shape.table
    # 在末列（col index 3）每行写入文本
    for row_idx in range(len(table.rows)):
        pptx_ops.set_cell(tbl_shape, row_idx, 3, f"末列第{row_idx}行")
    # 增加一列（第5列，index 4）
    pptx_ops.set_table_size(tbl_shape, rows=3, cols=5)
    assert len(table.columns) == 5
    # 新增列的每个单元格应为空
    for row_idx in range(len(table.rows)):
        new_cell_text = table.cell(row_idx, 4).text
        assert new_cell_text == "", (
            f"第{row_idx}行新列单元格应为空，但得到 {new_cell_text!r}"
        )
