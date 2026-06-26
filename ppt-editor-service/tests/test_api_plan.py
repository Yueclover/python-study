import io
from pptx import Presentation
from fastapi.testclient import TestClient
import app.main as main_mod
from app.main import app


def _setup(tmp_path):
    from app.storage import Storage
    main_mod.storage = Storage(str(tmp_path))


def test_apply_plan_fill_repeat_drop(tmp_path, basic_pptx_path):
    _setup(tmp_path)
    client = TestClient(app)
    with open(basic_pptx_path, "rb") as f:
        doc = client.post("/parse", files={"file": ("t.pptx", f, "application/octet-stream")}).json()
    doc_id = doc["doc_id"]

    plan = [
        {"kind": "fill", "slide_id": "s1", "fields": {"s1_sh1": "新封面标题"}},
        {"kind": "repeat", "slide_id": "s2", "items": [
            {"s2_sh1": "要点一"}, {"s2_sh1": "要点二"}, {"s2_sh1": "要点三"}]},
        {"kind": "drop", "slide_id": "s1"},
    ]
    resp = client.post("/apply_plan", json={"doc_id": doc_id, "plan": plan})
    assert resp.status_code == 200
    body = resp.json()
    assert body["rejected"] == []
    assert body["ops_count"] >= 5
    name = body["download_url"].split("/files/")[1]

    out = client.get(f"/files/{name}")
    prs = Presentation(io.BytesIO(out.content))
    titles = [s.shapes[0].text_frame.text for s in prs.slides]
    # 原 s1 已删；剩 原s2 + 3 张副本
    assert len(prs.slides) == 4
    assert "要点一" in titles and "要点三" in titles


def test_apply_plan_unknown_doc(tmp_path):
    _setup(tmp_path)
    client = TestClient(app)
    resp = client.post("/apply_plan", json={"doc_id": "deadbeef", "plan": []})
    assert resp.status_code == 404


def test_apply_plan_malformed_doc(tmp_path):
    _setup(tmp_path)
    client = TestClient(app)
    resp = client.post("/apply_plan", json={"doc_id": "../../etc", "plan": []})
    assert resp.status_code == 404
