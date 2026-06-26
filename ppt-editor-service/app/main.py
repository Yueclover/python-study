import os

from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.responses import FileResponse
from pptx import Presentation

from .storage import Storage, valid_doc_id
from .parser import parse_presentation
from .skeleton import build_skeleton
from .applier import apply_ops
from .plan import expand_plan
from .models import ApplyRequest, ApplyPlanRequest
from .validate import validate_html, repair_html

app = FastAPI(title="PPT Editor Service")
storage = Storage(os.environ.get("PPT_STORAGE", "./storage"))


@app.get("/health")
def health():
    return {"status": "ok"}


@app.post("/parse")
async def parse_endpoint(file: UploadFile = File(...)):
    data = await file.read()
    doc_id = storage.new_doc(data)
    try:
        prs = Presentation(storage.source_path(doc_id))
    except Exception:
        raise HTTPException(status_code=400, detail="无法解析 pptx 文件")
    result = parse_presentation(prs)
    result["doc_id"] = doc_id
    result["skeleton"] = build_skeleton(result)
    return result


@app.post("/apply")
def apply_endpoint(req: ApplyRequest):
    if not valid_doc_id(req.doc_id) or not storage.exists(req.doc_id):
        raise HTTPException(status_code=404, detail="doc_id 不存在")
    prs = Presentation(storage.source_path(req.doc_id))
    applied, rejected = apply_ops(prs, req.ops)
    out = storage.output_path(req.doc_id)
    prs.save(out)
    try:
        Presentation(out)  # 输出有效性校验
    except Exception:
        raise HTTPException(status_code=500, detail="生成的 pptx 校验失败")
    name = os.path.basename(out)
    return {"download_url": f"/files/{name}", "applied": applied, "rejected": rejected}


@app.post("/apply_plan")
def apply_plan_endpoint(req: ApplyPlanRequest):
    if not valid_doc_id(req.doc_id) or not storage.exists(req.doc_id):
        raise HTTPException(status_code=404, detail="doc_id 不存在")
    prs = Presentation(storage.source_path(req.doc_id))
    structure = parse_presentation(prs)
    ops, warnings = expand_plan(req.plan, structure)
    applied, rejected = apply_ops(prs, ops)
    out = storage.output_path(req.doc_id)
    prs.save(out)
    try:
        Presentation(out)
    except Exception:
        raise HTTPException(status_code=500, detail="生成的 pptx 校验失败")
    name = os.path.basename(out)
    return {"download_url": f"/files/{name}", "applied": applied,
            "rejected": rejected, "ops_count": len(ops), "warnings": warnings}


@app.post("/validate")
def validate_endpoint(req: dict):
    return validate_html(req.get("html", ""))


@app.post("/fix")
def fix_endpoint(req: dict):
    """输入 {"html": htmlStr}，等比缩放修复 section 内容溢出，返回 {"html": htmlStr}。"""
    return {"html": repair_html(req.get("html", ""))}


@app.get("/files/{name}")
def download(name: str):
    doc_id = name.split("-out")[0]
    if not valid_doc_id(doc_id):
        raise HTTPException(status_code=404, detail="文件不存在")
    path = storage.output_path(doc_id)
    if not (os.path.basename(path) == name and os.path.isfile(path)):
        raise HTTPException(status_code=404, detail="文件不存在")
    return FileResponse(
        path,
        filename=name,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
