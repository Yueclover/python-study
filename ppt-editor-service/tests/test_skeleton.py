from pptx import Presentation
from app.parser import parse_presentation
from app.skeleton import build_skeleton, _page_role


def test_skeleton_cover_and_content(basic_pptx_path):
    parsed = parse_presentation(Presentation(basic_pptx_path))
    sk = build_skeleton(parsed)
    assert len(sk["slides"]) == 2
    s1 = sk["slides"][0]
    assert s1["slide_id"] == "s1"
    assert s1["role"] == "cover"            # title + subtitle
    roles = {slot["role"] for slot in s1["slots"]}
    assert "title" in roles and "subtitle" in roles
    # 标题槽 shape_id 与解析一致
    title_slot = next(s for s in s1["slots"] if s["role"] == "title")
    assert title_slot["shape_id"] == "s1_sh1"
    assert sk["slides"][1]["role"] == "content"   # title + body


def test_skeleton_table_page(table_pptx_path):
    parsed = parse_presentation(Presentation(table_pptx_path))
    sk = build_skeleton(parsed)
    page = sk["slides"][0]
    assert page["role"] == "table"
    tslot = next(s for s in page["slots"] if s["role"] == "table")
    assert tslot["type"] == "table"
    assert tslot["rows"] == 3 and tslot["cols"] == 4


def test_slot_has_full_current_text_and_editable(basic_pptx_path):
    parsed = parse_presentation(Presentation(basic_pptx_path))
    sk = build_skeleton(parsed)
    title_slot = next(s for s in sk["slides"][0]["slots"] if s["role"] == "title")
    assert title_slot["current_text"] == "原标题"     # 全文，不截断
    assert title_slot["editable"] is True
    assert "sample" not in title_slot


def test_table_slot_not_editable(table_pptx_path):
    parsed = parse_presentation(Presentation(table_pptx_path))
    sk = build_skeleton(parsed)
    tslot = next(s for s in sk["slides"][0]["slots"] if s["role"] == "table")
    assert tslot["editable"] is False


def test_skeleton_exposes_layout_name(basic_pptx_path):
    parsed = parse_presentation(Presentation(basic_pptx_path))
    sk = build_skeleton(parsed)
    # python-pptx 默认模板：第一页版式名为 "Title Slide"
    assert sk["slides"][0]["layout_name"] == "Title Slide"


def test_page_role_table_wins():
    assert _page_role(["title", "table"], "Whatever", "", 1, 3) == "table"


def test_page_role_layout_cover():
    assert _page_role(["title"], "Title Slide", "", 1, 3) == "cover"
    assert _page_role(["title"], "封面页", "", 1, 3) == "cover"


def test_page_role_layout_toc():
    assert _page_role(["title", "body"], "目录", "", 1, 3) == "toc"
    assert _page_role(["title"], "Agenda", "", 1, 3) == "toc"


def test_page_role_layout_section():
    assert _page_role(["title"], "Section Header", "", 1, 3) == "section"
    assert _page_role(["title"], "节标题", "", 1, 3) == "section"


def test_page_role_layout_ending():
    assert _page_role(["title"], "结束页", "", 2, 3) == "ending"


def test_page_role_text_keyword_toc():
    assert _page_role(["title", "body"], "Blank", "目录 第一部分 第二部分", 1, 3) == "toc"


def test_page_role_text_keyword_ending():
    assert _page_role(["other"], "Blank", "谢谢观看", 2, 3) == "ending"


def test_page_role_text_keyword_section():
    assert _page_role(["title"], "Blank", "第一章 总览", 1, 4) == "section"


def test_page_role_index0_title_is_cover():
    assert _page_role(["title"], "Blank", "某标题", 0, 3) == "cover"


def test_page_role_last_sparse_is_ending():
    assert _page_role(["other"], "Blank", "再见", 2, 3) == "ending"


def test_page_role_composition_fallbacks():
    assert _page_role(["title", "subtitle"], "Blank", "x", 1, 3) == "cover"
    assert _page_role(["title", "body"], "Blank", "x", 1, 3) == "content"


def test_page_role_generic_default():
    assert _page_role(["other"], "Blank", "一些正文 内容很多 不止一个槽",
                      1, 3) == "generic"
