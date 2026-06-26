"""Demo 辅助:从 python-pptx 抽取背景 / 形状填充 / 描边 / 几何,供 HTML 还原。

边界(已知降级,均为 demo 取舍,不碰生产 parser):
- 主题色:解析 theme clrScheme;但忽略 lumMod/lumOff 明暗微调。
- 渐变:线性渐变转 CSS linear-gradient(角度近似);径向等回退为首个色标纯色。
- 几何:椭圆/圆角矩形/三角形可映射;FREEFORM 任意路径只能用包围盒近似。
"""
import math

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL
from pptx.oxml.ns import qn

# MSO_THEME_COLOR.name -> theme clrScheme 元素名
_THEME_NAME = {
    "DARK_1": "dk1", "LIGHT_1": "lt1", "DARK_2": "dk2", "LIGHT_2": "lt2",
    "TEXT_1": "dk1", "BACKGROUND_1": "lt1", "TEXT_2": "dk2", "BACKGROUND_2": "lt2",
    "ACCENT_1": "accent1", "ACCENT_2": "accent2", "ACCENT_3": "accent3",
    "ACCENT_4": "accent4", "ACCENT_5": "accent5", "ACCENT_6": "accent6",
    "HYPERLINK": "hlink", "FOLLOWED_HYPERLINK": "folHlink",
}
# XML schemeClr@val -> theme clrScheme 元素名
_SCHEME_REF = {
    "bg1": "lt1", "tx1": "dk1", "bg2": "lt2", "tx2": "dk2",
    "dk1": "dk1", "lt1": "lt1", "dk2": "dk2", "lt2": "lt2",
    "accent1": "accent1", "accent2": "accent2", "accent3": "accent3",
    "accent4": "accent4", "accent5": "accent5", "accent6": "accent6",
    "hlink": "hlink", "folHlink": "folHlink",
}


def load_theme(prs):
    """读 master 关联的 theme1.xml,返回 {dk1:'#RRGGBB', accent1:..., ...}。"""
    from lxml import etree
    out = {}
    try:
        master_part = prs.slide_masters[0].part
        theme_part = None
        for rel in master_part.rels.values():
            if "theme" in rel.reltype:
                theme_part = rel.target_part
                break
        if theme_part is None:
            return out
        # theme 是通用 Part,只有 .blob 原始字节,自己用 lxml 解析
        root = etree.fromstring(theme_part.blob)
        scheme = root.find(qn("a:themeElements") + "/" + qn("a:clrScheme"))
        if scheme is None:
            return out
        for child in scheme:
            key = child.tag.split("}")[-1]  # dk1 / lt1 / accent1 ...
            srgb = child.find(qn("a:srgbClr"))
            sysc = child.find(qn("a:sysClr"))
            if srgb is not None:
                out[key] = "#" + srgb.get("val")
            elif sysc is not None and sysc.get("lastClr"):
                out[key] = "#" + sysc.get("lastClr")
    except Exception:
        pass
    return out


def _hex_from_xml_color(el, theme):
    """el 是 <a:srgbClr> 或 <a:schemeClr>,返回 '#RRGGBB' 或 None。"""
    if el is None:
        return None
    tag = el.tag.split("}")[-1]
    if tag == "srgbClr":
        return "#" + el.get("val")
    if tag == "schemeClr":
        key = _SCHEME_REF.get(el.get("val"))
        return theme.get(key) if key else None
    return None


def _color_of(child_holder, theme):
    """在一个 XML 节点下找首个颜色子元素(srgbClr/schemeClr)。"""
    for tag in ("a:srgbClr", "a:schemeClr"):
        el = child_holder.find(qn(tag))
        if el is not None:
            return _hex_from_xml_color(el, theme)
    return None


def resolve_fore_color(color_format, theme):
    """FillFormat.fore_color -> '#hex',兼容 RGB 与 SCHEME 主题色。"""
    try:
        return "#" + str(color_format.rgb)
    except Exception:
        pass
    try:
        name = color_format.theme_color.name  # e.g. BACKGROUND_1
        key = _THEME_NAME.get(name)
        return theme.get(key) if key else None
    except Exception:
        return None


def _gradient_stops(spPr, theme):
    """<a:gradFill> -> ([(pos%,#hex), ...] 已排序, css_deg) 或 None。"""
    grad = spPr.find(qn("a:gradFill"))
    if grad is None:
        return None
    gs_lst = grad.find(qn("a:gsLst"))
    if gs_lst is None:
        return None
    stops = []
    for gs in gs_lst.findall(qn("a:gs")):
        pos = int(gs.get("pos", "0")) / 1000.0  # 0..100000 -> 0..100
        color = _color_of(gs, theme)
        if color:
            stops.append((pos, color))
    if not stops:
        return None
    stops.sort(key=lambda x: x[0])
    lin = grad.find(qn("a:lin"))
    if lin is not None:
        ang_deg = int(lin.get("ang", "0")) / 60000.0  # PPTX:从东向顺时针
        css_deg = (ang_deg + 90) % 360  # 近似映射到 CSS(0deg=向上)
    else:
        css_deg = 180  # 径向等不支持 → 默认从上到下
    return stops, css_deg


def _gradient_css(spPr, theme):
    """<a:gradFill> -> CSS linear-gradient(...) 或纯色回退,失败返回 None。"""
    res = _gradient_stops(spPr, theme)
    if not res:
        return None
    stops, css_deg = res
    if len(stops) == 1:
        return stops[0][1]  # 纯色回退
    body = ", ".join(f"{c} {round(p)}%" for p, c in stops)
    return f"linear-gradient({round(css_deg)}deg, {body})"


def _fill_css(shape, theme):
    """返回 background 的 CSS 值(纯色 hex 或 linear-gradient),无填充返回 None。"""
    spPr = getattr(shape, "_element", None)
    spPr = spPr.spPr if spPr is not None and hasattr(spPr, "spPr") else None
    try:
        ft = shape.fill.type
    except Exception:
        return None
    if ft == MSO_FILL.SOLID:
        return resolve_fore_color(shape.fill.fore_color, theme)
    if ft == MSO_FILL.GRADIENT and spPr is not None:
        return _gradient_css(spPr, theme)
    return None


def _line_css(shape, theme):
    """返回 border 的 CSS 值,无描边返回 None。"""
    try:
        line = shape.line
        color = resolve_fore_color(line.color, theme)
        if not color:
            return None
        w = line.width
        wpx = max(1, round(int(w) / 9525)) if w else 1
        return f"{wpx}px solid {color}"
    except Exception:
        return None


def _geom(shape):
    """oval / rounded / triangle / None。"""
    if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
        return None
    try:
        name = str(shape.auto_shape_type)
    except Exception:
        return None
    if "OVAL" in name:
        return "oval"
    if "ROUNDED_RECTANGLE" in name:
        return "rounded"
    if "TRIANGLE" in name:
        return "triangle"
    return None


def _pts(cmd):
    return [(int(p.get("x")), int(p.get("y"))) for p in cmd.findall(qn("a:pt"))]


def _custgeom_paths(spPr):
    """custGeom -> (viewbox_w, viewbox_h, [d_string, ...]) 或 None。

    支持 moveTo/lnTo/cubicBezTo/quadBezTo/close;arcTo 用折线采样近似。
    """
    cg = spPr.find(qn("a:custGeom"))
    if cg is None:
        return None
    path_lst = cg.find(qn("a:pathLst"))
    if path_lst is None:
        return None
    vb_w = vb_h = 0
    ds = []
    for path in path_lst.findall(qn("a:path")):
        vb_w = max(vb_w, int(path.get("w") or 0))
        vb_h = max(vb_h, int(path.get("h") or 0))
        d = []
        cur = (0, 0)
        for cmd in path:
            tag = cmd.tag.split("}")[-1]
            if tag == "moveTo":
                (x, y), = _pts(cmd)
                d.append(f"M{x} {y}")
                cur = (x, y)
            elif tag == "lnTo":
                (x, y), = _pts(cmd)
                d.append(f"L{x} {y}")
                cur = (x, y)
            elif tag == "cubicBezTo":
                (x1, y1), (x2, y2), (x, y) = _pts(cmd)
                d.append(f"C{x1} {y1} {x2} {y2} {x} {y}")
                cur = (x, y)
            elif tag == "quadBezTo":
                (x1, y1), (x, y) = _pts(cmd)
                d.append(f"Q{x1} {y1} {x} {y}")
                cur = (x, y)
            elif tag == "arcTo":
                cur = _arc_to_lines(cmd, cur, d)
            elif tag == "close":
                d.append("Z")
        if d:
            ds.append(" ".join(d))
    if not ds or vb_w == 0 or vb_h == 0:
        return None
    return vb_w, vb_h, ds


def _arc_to_lines(cmd, cur, d):
    """arcTo 折线采样(本 deck 用不到,留作兜底);返回新的 current 点。"""
    wr = int(cmd.get("wR") or 0)
    hr = int(cmd.get("hR") or 0)
    st = int(cmd.get("stAng") or 0) / 60000.0 * math.pi / 180
    sw = int(cmd.get("swAng") or 0) / 60000.0 * math.pi / 180
    cx = cur[0] - wr * math.cos(st)
    cy = cur[1] - hr * math.sin(st)
    steps = max(2, int(abs(sw) / (math.pi / 18)))
    x, y = cur
    for i in range(1, steps + 1):
        a = st + sw * i / steps
        x = cx + wr * math.cos(a)
        y = cy + hr * math.sin(a)
        d.append(f"L{x:.0f} {y:.0f}")
    return (x, y)


def freeform_svg(shape, theme, uid):
    """FREEFORM -> 内嵌 <svg> 字符串(含纯色/渐变填充),失败返回 None。"""
    spPr = getattr(getattr(shape, "_element", None), "spPr", None)
    if spPr is None:
        return None
    geom = _custgeom_paths(spPr)
    if geom is None:
        return None
    vb_w, vb_h, ds = geom

    defs = ""
    fill = "#cccccc"
    try:
        ft = shape.fill.type
    except Exception:
        ft = None
    if ft == MSO_FILL.SOLID:
        fill = resolve_fore_color(shape.fill.fore_color, theme) or fill
    elif ft == MSO_FILL.GRADIENT:
        res = _gradient_stops(spPr, theme)
        if res:
            stops, css_deg = res
            if len(stops) == 1:
                fill = stops[0][1]
            else:
                gid = "g" + uid.replace("_", "")
                rot = round(css_deg - 90)  # SVG 默认 0deg=向右
                stop_tags = "".join(
                    f"<stop offset='{round(p)}%' stop-color='{c}'/>"
                    for p, c in stops)
                defs = (
                    f"<defs><linearGradient id='{gid}' "
                    f"gradientUnits='objectBoundingBox' "
                    f"gradientTransform='rotate({rot} 0.5 0.5)'>"
                    f"{stop_tags}</linearGradient></defs>")
                fill = f"url(#{gid})"

    paths = "".join(f"<path d='{d}' fill='{fill}'/>" for d in ds)
    return (
        f"<svg viewBox='0 0 {vb_w} {vb_h}' preserveAspectRatio='none' "
        f"style='width:100%;height:100%;display:block;overflow:visible'>"
        f"{defs}{paths}</svg>")


def _flip_rot(shape):
    """读 xfrm 的 rot/flipH/flipV,返回 CSS transform 片段(可能为空)。"""
    spPr = getattr(getattr(shape, "_element", None), "spPr", None)
    if spPr is None:
        return ""
    xfrm = spPr.find(qn("a:xfrm"))
    if xfrm is None:
        return ""
    parts = []
    rot = int(xfrm.get("rot") or 0) / 60000.0
    if rot:
        parts.append(f"rotate({round(rot, 2)}deg)")
    sx = -1 if xfrm.get("flipH") == "1" else 1
    sy = -1 if xfrm.get("flipV") == "1" else 1
    if sx != 1 or sy != 1:
        parts.append(f"scale({sx},{sy})")
    return ("transform:" + " ".join(parts) + ";") if parts else ""


def extract_visuals(prs, theme):
    """{shape_id: {bg, border, geom, is_freeform}} —— 与 parser 同序。"""
    from app.ids import shape_id_for
    out = {}
    for si, slide in enumerate(prs.slides):
        for oi, shp in enumerate(slide.shapes):
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                continue  # 图片单独走 extract_images
            sid = shape_id_for(si, oi)
            is_ff = shp.shape_type == MSO_SHAPE_TYPE.FREEFORM
            info = {
                "bg": _fill_css(shp, theme),
                "border": _line_css(shp, theme),
                "geom": _geom(shp),
                "is_freeform": is_ff,
                "svg": freeform_svg(shp, theme, sid) if is_ff else None,
                "transform": _flip_rot(shp),
            }
            if any(v for k, v in info.items() if k != "transform"):
                out[sid] = info
    return out


def slide_bg_css(slide, prs, theme):
    """页面背景:纯色 → 渐变 → None(交给 CSS 默认白)。"""
    try:
        bg = slide.background
        ft = bg.fill.type
        if ft == MSO_FILL.SOLID:
            return resolve_fore_color(bg.fill.fore_color, theme)
        if ft == MSO_FILL.GRADIENT:
            spPr = bg.fill._xPr if hasattr(bg.fill, "_xPr") else None
            if spPr is not None:
                return _gradient_css(spPr, theme)
    except Exception:
        pass
    return None
